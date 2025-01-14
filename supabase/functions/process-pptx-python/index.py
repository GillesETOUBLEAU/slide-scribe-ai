import os
import json
from datetime import datetime
from pptx import Presentation
from supabase import create_client, Client
from typing import Dict, Any

# Initialize Supabase client
supabase: Client = create_client(
    os.environ.get("SUPABASE_URL"),
    os.environ.get("SUPABASE_SERVICE_ROLE_KEY")
)

def cors_headers():
    return {
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Headers": "authorization, x-client-info, apikey, content-type",
    }

def process_pptx(file_path: str) -> Dict[str, Any]:
    print(f"Processing PPTX file: {file_path}")
    prs = Presentation(file_path)
    slides = []
    
    for idx, slide in enumerate(prs.slides, 1):
        content = []
        notes = []
        
        # Extract slide content
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content.append(shape.text.strip())
                
        # Extract notes if they exist
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes.append(slide.notes_slide.notes_text_frame.text.strip())
        
        slides.append({
            "index": idx,
            "title": content[0] if content else f"Slide {idx}",
            "content": content[1:] if content else [],
            "notes": notes,
            "shapes": []  # We'll keep this for compatibility
        })
    
    return {
        "metadata": {
            "processedAt": datetime.now().isoformat(),
            "filename": os.path.basename(file_path),
            "slideCount": len(slides)
        },
        "slides": slides
    }

async def handle_request(req):
    # Handle CORS preflight requests
    if req.method == "OPTIONS":
        return new Response(None, { "headers": cors_headers() })
    
    try:
        body = await req.json()
        file_id = body.get("fileId")
        file_path = body.get("filePath")
        
        if not file_id or not file_path:
            raise ValueError("Missing required parameters: fileId or filePath")

        print(f"Starting processing for file: {file_id}")
        
        # Download file from storage
        response = await supabase.storage.from_("pptx_files").download(file_path)
        if not response.data:
            raise Exception("Failed to download file from storage")
            
        # Save temporarily and process
        temp_path = f"/tmp/{os.path.basename(file_path)}"
        with open(temp_path, "wb") as f:
            f.write(response.data)
            
        # Process the file
        result = process_pptx(temp_path)
        
        # Clean up temp file
        os.remove(temp_path)
        
        # Generate output paths
        json_path = file_path.replace(".pptx", ".json")
        markdown_path = file_path.replace(".pptx", ".md")
        
        # Upload JSON result
        json_content = json.dumps(result)
        await supabase.storage.from_("pptx_files").upload(
            json_path,
            json_content,
            {"contentType": "application/json"}
        )
        
        # Generate and upload markdown
        markdown_content = generate_markdown(result)
        await supabase.storage.from_("pptx_files").upload(
            markdown_path,
            markdown_content,
            {"contentType": "text/markdown"}
        )
        
        # Update database record
        await supabase.table("file_conversions").update({
            "status": "completed",
            "json_path": json_path,
            "markdown_path": markdown_path
        }).eq("id", file_id).execute()
        
        return new Response(
            json.dumps({"success": True}),
            headers={**cors_headers(), "Content-Type": "application/json"}
        )
        
    except Exception as e:
        print(f"Error processing file: {str(e)}")
        if file_id:
            await supabase.table("file_conversions").update({
                "status": "error",
                "error_message": str(e)
            }).eq("id", file_id).execute()
        
        return new Response(
            json.dumps({"error": str(e)}),
            status=500,
            headers={**cors_headers(), "Content-Type": "application/json"}
        )

def generate_markdown(data: Dict[str, Any]) -> str:
    markdown = f"# {data['metadata']['filename']}\n\n"
    markdown += f"Processed at: {data['metadata']['processedAt']}\n"
    markdown += f"Total Slides: {data['metadata']['slideCount']}\n\n"
    
    for slide in data['slides']:
        markdown += f"## {slide['title']}\n\n"
        
        if slide['content']:
            markdown += "### Content\n\n"
            for content in slide['content']:
                markdown += f"- {content}\n"
            markdown += "\n"
            
        if slide['notes']:
            markdown += "### Notes\n\n"
            for note in slide['notes']:
                markdown += f"> {note}\n"
            markdown += "\n"
    
    return markdown

serve(handle_request)